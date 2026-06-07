"""
document_service.py — Text extraction for Knowledge Vault (Day 7).

extract_text(file_path, file_type) -> str

Supports PDF (pypdf), DOCX (python-docx), PPTX (python-pptx), TXT.
Falls back gracefully if a package is not installed.
"""

import logging

logger = logging.getLogger(__name__)

MAX_CHARS = 40_000   # ~10 000 tokens — reasonable AI context limit


def extract_text(file_path: str, file_type: str) -> str:
    """
    Extract plain text from a document at *file_path*.
    Returns empty string on failure (never raises).
    """
    try:
        ext = file_type.lower().strip('.')
        if ext == 'txt':
            return _extract_txt(file_path)
        if ext == 'pdf':
            return _extract_pdf(file_path)
        if ext == 'docx':
            return _extract_docx(file_path)
        if ext == 'pptx':
            return _extract_pptx(file_path)
        logger.warning('[Document] Unknown file type: %s', ext)
        return ''
    except Exception as exc:
        logger.error('[Document] Extraction failed (%s): %s', file_path, exc)
        return ''


# ---------------------------------------------------------------------------
# Internal extractors
# ---------------------------------------------------------------------------

def _extract_txt(path: str) -> str:
    with open(path, 'r', encoding='utf-8', errors='ignore') as fh:
        return fh.read()[:MAX_CHARS]


def _extract_pdf(path: str) -> str:
    try:
        import pypdf
    except ImportError:
        logger.warning('[Document] pypdf not installed — run: pip install pypdf')
        return ''

    reader = pypdf.PdfReader(path)
    pages  = [page.extract_text() or '' for page in reader.pages]
    return '\n'.join(pages)[:MAX_CHARS]


def _extract_docx(path: str) -> str:
    try:
        import docx
    except ImportError:
        logger.warning('[Document] python-docx not installed — run: pip install python-docx')
        return ''

    doc   = docx.Document(path)
    paras = [p.text for p in doc.paragraphs if p.text.strip()]
    return '\n'.join(paras)[:MAX_CHARS]


def _extract_pptx(path: str) -> str:
    try:
        from pptx import Presentation
    except ImportError:
        logger.warning('[Document] python-pptx not installed — run: pip install python-pptx')
        return ''

    prs   = Presentation(path)
    texts = []
    for slide in prs.slides:
        for shape in slide.shapes:
            if hasattr(shape, 'text') and shape.text.strip():
                texts.append(shape.text)
    return '\n'.join(texts)[:MAX_CHARS]
